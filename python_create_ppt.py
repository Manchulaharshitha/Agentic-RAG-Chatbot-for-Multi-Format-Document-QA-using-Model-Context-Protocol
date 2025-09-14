from pptx import Presentation
# Create a new presentation
prs = Presentation()
slides_content = [
    ("Agentic RAG Chatbot for Multi-Format Document QA",
     "Using FAISS, GPT, and Streamlit for intelligent document question answering\n\nYour Name | Date"),
    ("Introduction",
     "• AI chatbot to answer questions from PDFs, DOCX, TXT.\n• Uses embeddings, FAISS, GPT for accurate responses.\n• Interactive Streamlit interface."),
    ("Problem Statement",
     "• Users struggle to find info from large documents.\n• Manual searching is time-consuming.\n• Need for flexible multi-format solution."),
    ("RAG Overview",
     "• Combines semantic search with GPT.\n• Retrieves relevant context.\n• Generates human-like answers."),
    ("MCP (Model Context Protocol)",
     "• Coordinates ingestion, retrieval, and LLM agents.\n• Ensures smooth workflow."),
    ("System Architecture",
     "• Ingestion Agent: Reads files.\n• Retrieval Agent: Embeddings & FAISS.\n• LLM Agent: GPT responses.\n• Streamlit UI for interaction."),
    ("Key Technologies",
     "• Python 3.x\n• Streamlit\n• FAISS\n• Sentence Transformers\n• OpenAI GPT\n• PyPDF2 / python-docx"),
    ("Code Walkthrough",
     "• ingestion_agent.py – Parse files\n• retrieval_agent.py – Embeddings & search\n• llm_response_agent.py – GPT responses\n• mcp.py – Agent coordination\n• app.py – Streamlit UI"),
    ("Demo / Workflow",
     "1. Upload document → parse\n2. Ask question → search context\n3. GPT generates response → displayed"),
    ("Challenges and Solutions",
     "Challenges:\n• Multi-format files\n• Accurate retrieval\n• Maintaining context\nSolutions:\n• Robust libraries\n• FAISS & embeddings\n• MCP"),
    ("Future Work",
     "• Add more formats\n• Domain-specific embeddings\n• Cloud deployment\n• Authentication & privacy\n• Multilingual support"),
    ("Conclusion & Contact",
     "• Efficient document QA chatbot\n• Streamlined access to information\n• Real-world AI application\n\nContact:\nEmail: manchulaharshitha12@gmail.com\nGitHub: https://github.com/Manchulaharshitha/rag-chatbot")
]
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content layout
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
prs.save("rag_chatbot_presentation.pptx")
print("Presentation created successfully!")
